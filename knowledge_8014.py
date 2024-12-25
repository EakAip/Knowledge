# 接口：8014

# PPT知识图谱

# 支持：pptx,ppt,pdf

# 本地模型 切分PPT处理

# 以1200字为上限 动态切分

from flask import Flask,request,jsonify
from openai import OpenAI
from pptx import Presentation
import pdfplumber
import os
import re
import json
import time
import logging
import threading
import subprocess
from urllib.parse import urlparse
import requests # 用于发送回调请求

app = Flask(__name__)

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def convert_ppt_to_pptx(input_file, output_dir):
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    # 调用libreoffice命令进行转换
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pptx",
        input_file,
        "--outdir",
        output_dir
    ]
    subprocess.run(command, check=True)
    print(f"Converted {input_file} to .pptx format in {output_dir}")


# 提取PDF
def extract_text_from_pdf(file_path):
    all_text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                all_text += text + "\n"
                
    return all_text

# 提取PPTX
def extract_text_from_pptx(file_path):
    all_text = ""

    # 如果文件是 .ppt 格式，则先转换为 .pptx
    if file_path.endswith('.ppt'):
        print(f"Detected .ppt file, converting {file_path} to .pptx...")
        output_dir = os.path.dirname(file_path)  # 使用相同目录
        try:
            convert_ppt_to_pptx(file_path, output_dir)
            # 更新文件路径为 .pptx
            file_path = file_path.replace('.ppt', '.pptx')
        except subprocess.CalledProcessError as e:
            print(f"Error during conversion: {e}")
            raise ValueError("Failed to convert .ppt to .pptx. Please check the file or LibreOffice installation.")
    
    # 加载 .pptx 文件提取文本
    try:
        presentation = Presentation(file_path)
        # 遍历每个幻灯片
        for slide in presentation.slides:
            # 遍历每个形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading .pptx file: {e}")
        raise ValueError("Failed to read the .pptx file. Ensure the file is valid.")
    
    # # 将all_text保存到当前目录下
    # with open("res.txt", "w", encoding="utf-8") as f:
    #     f.write(all_text)

    return all_text


# 定义提取知识点解析函数
def parse_text_to_structure(text):
    lines = text.split('\n')
    lines = [line.lstrip() for line in lines]  # 去除每行前面的空白字符
    structure = {"nodes": [], "message": "成功", "code": 0}
    levels = {}

    for line in lines:
        if not line.strip():
            continue

        chapter_match = re.match(r'第\s*([一二三四五六七八九十\d]+)\s*章[:：]?\s*(.*)', line)
        level_matches = re.findall(r'(\d+(?:\.\d+)*)\s+(.*)', line)

        if chapter_match:
            chapter_number, chapter_title = chapter_match.groups()
            chapter_title = chapter_title.strip()
            current_node = {"level": 1, "name": chapter_title, "child": [], "postponement": "", "definition": ""}
            structure['nodes'].append(current_node)
            levels[1] = current_node
        elif level_matches:
            for match in level_matches:
                level_str, title = match
                title = title.strip()
                level = len(level_str.split('.'))
                current_node = {"level": level, "name": title, "child": [], "postponement": "", "definition": ""}
                if level > 1 and (level - 1) in levels:
                    levels[level - 1]['child'].append(current_node)
                levels[level] = current_node
        else:
            print(f"无法解析的行,已经跳过：{line}")
            continue

    return structure

# 回调
def send_callback(rpId, authorization_token, callback_url, data,file_name):
    try:
        
        # 向回调接口发送内容
        callback_data = {
            "retJsonStr":data,
            "rpId":rpId,
        }
        
        headers = {
            'AuthorizationForPlatform': authorization_token
        }
        
        # 发送数据到回调接口
        response = requests.post(callback_url, json=callback_data,headers=headers)
        logging.info(f"[{file_name}]发送数据到回调接口")

        # 检查回调接口响应
        if response.status_code == 200:
            logging.info(f"[{file_name}]回调接口响应成功")
        else:
            logging.info(f"[{file_name}]回调接口响应失败")


    
    except requests.exceptions.RequestException as e:
        logging.info(f"[{file_name}]发送回调失败: {e}")


def api1(query):
    # API 基本信息和认证设置
    API_URL = "http://188.18.18.106:5001/v1/completion-messages"
    API_KEY = "app-Ord2AUBiP8ME4jBhQ5iRW5Wq"  # "生成目录" 工作流API密钥
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "inputs": {"query": query},
        "response_mode": "streaming",  # 设置为流式模式
        "user": "abc-123"
    }
    
    result = ""
    
    try:
        # 发起请求并获取响应
        with requests.post(API_URL, headers=headers, json=payload, stream=True, timeout=120) as response:
            if response.status_code != 200:
                logging.info(f"请求失败，状态码：{response.status_code}")
                try:
                    error_data = response.json()
                    logging.info(f"错误信息：{error_data}")
                except ValueError:
                    logging.info("响应内容不是有效的 JSON，无法获取错误信息")
                return result

            # 逐行读取响应内容
            for line in response.iter_lines():
                if line:
                    decoded_line = line.decode('utf-8').strip()
                    if decoded_line.startswith("data: "):
                        json_str = decoded_line[6:]  # 移除 "data: " 前缀
                        try:
                            data = json.loads(json_str)
                        except json.JSONDecodeError:
                            print(f"无法解析 JSON: {json_str}")
                            continue

                        # 处理不同的事件类型
                        event = data.get("event", "message")  # 默认为 'message' 类型
                        if event == "message":
                            answer = data.get("answer", "")
                            # 使用 end='' 和 flush=True 使输出不换行
                            print(answer, end='', flush=True)
                            result += answer
                        elif event == "tts_message":
                            audio = data.get("audio", "")
                            print(f"TTS 音频块: {audio}")
                        elif event == "tts_message_end":
                            # print("TTS 音频流结束")
                            pass
                        elif event == "message_end":
                            print("消息流结束")
                        elif event == "message_replace":
                            replaced_answer = data.get("answer", "")
                            print(f"替换后的回答: {replaced_answer}")
                        elif event == "error":
                            status = data.get("status", "")
                            code = data.get("code", "")
                            message = data.get("message", "")
                            print(f"错误事件 - 状态码: {status}, 错误码: {code}, 消息: {message}")
                        elif event == "ping":
                            print("收到 ping 事件，保持连接存活")
                        else:
                            print(f"未知事件类型: {event}")
    except requests.RequestException as e:
        logging.error(f"请求异常: {e}")
    
    return result

def api2(query,text):
    # API 基本信息和认证设置
    API_URL = "http://188.18.18.106:5001/v1/completion-messages"
    API_KEY = "app-leAyLRxGq5JQXTo80rcLcwNa"  # "生成目录" 工作流API密钥
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "inputs": {
            "query": query,
            "text": text,
            },
        "response_mode": "streaming",  # 设置为流式模式
        "user": "abc-123"
    }
    
    result = ""
    
    try:
        # 发起请求并获取响应
        with requests.post(API_URL, headers=headers, json=payload, stream=True, timeout=120) as response:
            if response.status_code != 200:
                logging.info(f"请求失败，状态码：{response.status_code}")
                try:
                    error_data = response.json()
                    logging.info(f"错误信息：{error_data}")
                except ValueError:
                    logging.info("响应内容不是有效的 JSON，无法获取错误信息")
                return result

            # 逐行读取响应内容
            for line in response.iter_lines():
                if line:
                    decoded_line = line.decode('utf-8').strip()
                    if decoded_line.startswith("data: "):
                        json_str = decoded_line[6:]  # 移除 "data: " 前缀
                        try:
                            data = json.loads(json_str)
                        except json.JSONDecodeError:
                            print(f"无法解析 JSON: {json_str}")
                            continue

                        # 处理不同的事件类型
                        event = data.get("event", "message")  # 默认为 'message' 类型
                        if event == "message":
                            answer = data.get("answer", "")
                            # 使用 end='' 和 flush=True 使输出不换行
                            print(answer, end='', flush=True)
                            result += answer
                        elif event == "tts_message":
                            audio = data.get("audio", "")
                            print(f"TTS 音频块: {audio}")
                        elif event == "tts_message_end":
                            # print("TTS 音频流结束")
                            pass
                        elif event == "message_end":
                            print("消息流结束")
                        elif event == "message_replace":
                            replaced_answer = data.get("answer", "")
                            print(f"替换后的回答: {replaced_answer}")
                        elif event == "error":
                            status = data.get("status", "")
                            code = data.get("code", "")
                            message = data.get("message", "")
                            print(f"错误事件 - 状态码: {status}, 错误码: {code}, 消息: {message}")
                        elif event == "ping":
                            print("收到 ping 事件，保持连接存活")
                        else:
                            print(f"未知事件类型: {event}")
    except requests.RequestException as e:
        logging.error(f"请求异常: {e}")
    
    return result



def process_file_and_callback(rpId, authorization_token, callback_url, file_path, file_name, modeltype):
    # 提取文本
    if file_path.endswith('.pdf'):
        all_text = extract_text_from_pdf(file_path)
        
    elif file_path.endswith('.pptx') or file_path.endswith('.ppt'):
        all_text = extract_text_from_pptx(file_path)
        
    else:
        logging.info(f"[{file_name}]不支持的文件类型")
        send_callback(callback_url, {"code": 1, "msg": "不支持的文件类型"})   # 这里是否有隐患？
        return

    # 处理all_text内容
    content = all_text.replace('\n',' ')
    # 剔除空格 
    content = content.replace('  ','')
    # 剔除"--"字符
    content = content.replace('--','')
    # 剔除"…"字符
    content = content.replace('…','')
    # 剔除". ."字符
    content = content.replace('. .','')
    
    # 判断字数是否超过15000字，不超过则只调用工作流接口1
    if len(content) < 15000:
        logging.info(f"[{file_name}] 文件字数小于15000字，不切分")
        res1 = content
        # 调用工作流接口1
    
        try:
            result = api1(res1)
            
            print(result)
            
            filename = file_name.split('.')[0]          # 不要file_name后缀
            # 将目录写入文件保存备份，将file_name.txt文件清空
            with open(f'data/catalog/{filename}目录.txt','w',encoding='utf-8') as f:
                # 先清空
                f.truncate()
                logging.info(f"[{file_name}] data/catalog/{filename}目录.txt——————>>>清空成功！")
                # 写入目录
                f.write(result)
                logging.info(f"[{file_name}]目录已经保存到: data/catalog/{filename}目录.txt")
            

            # 解析结构
            structure = parse_text_to_structure(result)
            
            print(structure)

            # 发送回调
            send_callback(rpId, authorization_token,callback_url, structure,file_name)
        except Exception as e:
            logging.error(f"[{file_name}] 处理失败: {e}")
        
    else:
        
        logging.info(f"[{file_name}]文件字数为{len(content)}，开始切分,切分为{len(content)//12000+1}份")
    
        split_size = 12000
        parts = []
        
        '''
            00000~~12000
            11500~~23500
            23000~~35000
            34500~~46500
            46000~~58000
        '''
        
        for i in range(0,len(content),split_size-500):   # 每部分向前多参考500字        
            parts.append(content[i:i+split_size])
            
        # 保存切分内容到文件res1.txt,res2.txt,res3.txt...
        # for i,part in enumerate(parts,start=1):
        #     with open(f"data/res{i}.txt",'w',encoding='utf-8') as f:
        #         f.truncate()
        #         f.write(part)
                
        try:
            # 调用工作流接口1
            logging.info(f"[{file_name}]调用工作流接口1")
            result1 = api1(parts[0])
            
            result = result1
            for i in range(1,len(parts)):
                # 调用工作流接口2
                logging.info(f"[{file_name}]调用工作流接口2,第{i+1}部分")
                result = api2(parts[i],result)
            print(result)
            
            filename = file_name.split('.')[0]          # 不要file_name后缀
            # 将目录写入文件保存备份，将file_name.txt文件清空
            with open(f'data/catalog/{filename}目录.txt','w',encoding='utf-8') as f:
                # 先清空
                f.truncate()
                logging.info(f"[{file_name}] data/catalog/{filename}目录.txt——————>>>清空成功！")
                # 写入目录
                f.write(result)
                logging.info(f"[{file_name}]目录已经保存到: data/catalog/{filename}目录.txt")

            # 解析结构
            structure = parse_text_to_structure(result)
            
            # print(structure)

            # 发送回调
            send_callback(rpId, authorization_token, callback_url, structure,file_name)

        except Exception as e:
            logging.error(f"[{file_name}] 处理失败: {e}")
                



@app.route('/knowledge', methods=['POST'])
def receive_knowledge():
    # 检查必填字段
    required_fields = ["rpId", "AuthorizationForPlatform", "callBackUrl"]

    for field in required_fields:
        if field not in request.form:
            return jsonify({"code": 5, 'msg': f"{field}字段缺失"}), 200

    rpId = request.form.get('rpId')
    authorization_token = request.form.get('AuthorizationForPlatform')
    callback_url = request.form.get('callBackUrl')
    logging.info(f"回调接口:{callback_url}")

    # 验证回调URL是否有效
    parsed_url = urlparse(callback_url)
    if not parsed_url.scheme or not parsed_url.netloc:
        return jsonify({"code": 5, 'msg': "回调URL无效"}), 200

    # 接收文件对象
    if 'text' not in request.files:
        return jsonify({"code": 5, 'msg': "缺少文件"}), 200

    file = request.files['text']
    if file.filename == '':
        return jsonify({"code": 5, 'msg': "未选择文件"}), 200

    # 检查文件扩展名
    allowed_extensions = {'pdf', 'pptx', 'ppt'}
    file_ext = os.path.splitext(file.filename)[1].lower().replace('.', '')
    if file_ext not in allowed_extensions:
        return jsonify({"code": 5, 'msg': "仅支持上传PDF或PPT文件"}), 200

    # 获取文件名字
    file_name = file.filename
    # 保存文件到服务器路径
    save_dir = 'data/ppt'
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, file_name)
    file.save(file_path)
    logging.info(f"文件已保存到 {file_path}")

    modeltype = request.form.get('modeltype', '1')  # 默认modeltype为'1'

    # 启动新线程处理文件并发送回调
    threading.Thread(target=process_file_and_callback, args=(rpId, authorization_token, callback_url, file_path, file_name, modeltype)).start()

    # 立即返回响应
    return jsonify({"code": 0, "rpId": rpId, "msg": "请求已接收，正在处理中......"}), 200

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8014)
    
